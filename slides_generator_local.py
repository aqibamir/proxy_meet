import os
import re
import json
import tempfile
from typing import List, Dict, Tuple
from uuid import uuid4
import logging

import openai
from openai import OpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def read_pdf(pdf_path: str) -> str:
    """Return raw text extracted from a PDF file."""
    if not pdf_path or not os.path.isfile(pdf_path):
        logger.warning(f"PDF file {pdf_path} not found or invalid.")
        return ""
    try:
        reader = PdfReader(pdf_path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        logger.info(f"Successfully extracted text from PDF: {pdf_path}")
        return text
    except Exception as e:
        logger.error(f"Error reading PDF {pdf_path}: {e}")
        return ""

def ask_llm(system: str, user: str, model: str = "gpt-4") -> str:
    """Send system/user prompts to OpenAI and return the text."""
    try:
        client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))
        if not client.api_key:
            raise ValueError("OPENAI_API_KEY environment variable not set")
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": user},
            ],
            max_tokens=2000,
            temperature=0.5,
        )
        content = response.choices[0].message.content
        logger.info("Successfully received LLM response")
        return content
    except Exception as e:
        logger.error(f"Error in LLM request: {e}")
        return ""

def parse_llm_to_slides(raw: str) -> List[Dict]:
    """Convert LLM plain-text response into structured slide objects."""
    slides: List[Dict] = []
    chunks = re.split(r"---|\n{2,}", raw.strip())
    logger.info(f"Parsed {len(chunks)} chunks from LLM response")
    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue

        title_match = re.search(r"Slide \d+ Title:\s*(.+)", chunk, re.I)
        type_match = re.search(r"Slide Type:\s*(.+)", chunk, re.I)
        body_match = re.search(r"Body:\s*([\s\S]+?)(?:\nTalking Points:|$)", chunk, re.I)
        talk_match = re.search(r"Talking Points:\s*([\s\S]+?)$", chunk, re.I)

        title = title_match.group(1).strip() if title_match else "Untitled"
        slide_type = type_match.group(1).strip().lower() if type_match else "content"
        body = body_match.group(1).strip() if body_match else ""
        talking = talk_match.group(1).strip() if talk_match else ""

        slides.append({
            "slide_number": len(slides) + 1,
            "slide_type": slide_type,
            "slide_title": title,
            "slide_content": body,
            "talking_points": talking,
        })
    logger.info(f"Generated {len(slides)} slides from LLM response")
    return slides

def build_pptx(slides: List[Dict], output_dir: str) -> str:
    """Create presentation.pptx in output_dir and return path."""
    try:
        prs = Presentation()
        
        # Set slide dimensions for widescreen (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_info in slides:
            slide_type = slide_info.get("slide_type", "content")
            title_text = slide_info["slide_title"]
            body_text = slide_info["slide_content"]
            layout_index = 1 if slide_type == "content" else 0  # 0: title-only, 1: title+content
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])

            # Style title
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_shape.text_frame.paragraphs[0].font.size = Pt(36)
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            if slide_type == "content" and body_text:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for line in body_text.splitlines():
                    line = line.strip("•-— ")
                    if not line:
                        continue
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0 if not line.startswith("  ") else 1  # Support nested bullets
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
                    p.space_before = Pt(10)
                    p.space_after = Pt(10)

            # Add image placeholder for visual slides
            if slide_type == "visual":
                placeholder = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(4), Inches(3))
                tf = placeholder.text_frame
                tf.text = "Image Placeholder"
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        pptx_path = os.path.join(output_dir, "presentation.pptx")
        prs.save(pptx_path)
        if os.path.exists(pptx_path):
            logger.info(f"PPTX saved successfully at {os.path.abspath(pptx_path)}")
        else:
            logger.error(f"PPTX file not found after saving: {pptx_path}")
            return ""
        return pptx_path
    except Exception as e:
        logger.error(f"Error saving PPTX: {e}")
        return ""

def build_json(slides: List[Dict], output_dir: str) -> str:
    """Create presentation_script.json in output_dir and return path."""
    try:
        json_path = os.path.join(output_dir, "presentation_script.json")
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(slides, f, indent=2, ensure_ascii=False)
        if os.path.exists(json_path):
            logger.info(f"JSON saved successfully at {os.path.abspath(json_path)}")
        else:
            logger.error(f"JSON file not found after saving: {json_path}")
            return ""
        return json_path
    except Exception as e:
        logger.error(f"Error saving JSON: {e}")
        return ""

def generate_local(prompt: str, pdf_path: str, output_dir: str) -> Tuple[str, str]:
    """Main entry point. Returns (pptx_file_path, json_file_path)."""
    try:
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            logger.info(f"Created output directory: {output_dir}")

        pdf_text = read_pdf(pdf_path)

        system_prompt = (
            "You are an expert presentation designer creating engaging, professional slides. "
            "Generate 3-5 slides based on the user prompt and PDF content. "
            "For each slide, output EXACTLY:\n"
            "Slide <#> Title: <clear, concise, and engaging title>\n"
            "Slide Type: <content|title|visual|quote>\n"
            "Body: <3-5 concise bullet points for 'content' slides, a single quote for 'quote' slides, or a description of a visual for 'visual' slides>\n"
            "Talking Points: <natural, presenter-friendly script (100-150 words)>\n"
            "Separate slides with '---'. Vary slide types for visual interest."
        )
        user_prompt = f"Prompt: {prompt}\n\nPDF Content:\n{pdf_text[:10000]}"  # Truncate to avoid token limits
        raw_content = ask_llm(system_prompt, user_prompt)
        logger.info(f"LLM Response:\n{raw_content}")

        slides = parse_llm_to_slides(raw_content)
        logger.info(f"Generated {len(slides)} slides")
        if not slides:
            logger.warning("No slides generated from LLM response.")
            return "", ""

        pptx_file = build_pptx(slides, output_dir)
        json_file = build_json(slides, output_dir)
        return pptx_file, json_file
    except Exception as e:
        logger.error(f"Error in generate_local: {e}")
        return "", ""

if __name__ == "__main__":
    try:
        prompt = "Create a presentation about the benefits of renewable energy."
        pdf_path = "renewable_energy.pdf"
        output_dir = "temp_files"
        pptx, json_file = generate_local(prompt, pdf_path, output_dir)
        if pptx and json_file:
            logger.info(f"Generated: {os.path.abspath(pptx)}, {os.path.abspath(json_file)}")
        else:
            logger.error("Failed to generate one or both files")
    except Exception as e:
        logger.error(f"Error in main execution: {e}")